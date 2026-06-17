"""
build_guide_deck.py  –  Clean user guide deck, no overlapping shapes.
Run:  python build_guide_deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).parent / "FIFA_Inspirational_App_Guide.pptx"

GREEN  = RGBColor(0x00, 0x6B, 0x3C)
GOLD   = RGBColor(0xD4, 0xAF, 0x37)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NAVY   = RGBColor(0x0D, 0x2A, 0x4A)
ORANGE = RGBColor(0xFF, 0x6B, 0x00)
BLUE   = RGBColor(0x1E, 0x90, 0xFF)
GREY   = RGBColor(0x99, 0x99, 0x99)
BLACK  = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Emu(9144000)   # 10 in widescreen
prs.slide_height = Emu(5143500)   # 7.5 in
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=DARK):
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


def box(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=12, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT):
    b  = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color


def footer(slide):
    txt(slide, "FIFA 2026 Inspirational Stories  ·  User Guide",
        Inches(0.3), H - Inches(0.35), W - Inches(0.6), Inches(0.3),
        size=9, color=GREY, align=PP_ALIGN.CENTER)


def title_bar(slide, title, subtitle=""):
    box(slide, 0, 0, W, Inches(1.05), GREEN)
    txt(slide, title,
        Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.55),
        size=22, bold=True, color=GOLD)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.3), Inches(0.62), W - Inches(0.6), Inches(0.38),
            size=11, italic=True, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1  Cover
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, GREEN)
box(s, 0, 0, Inches(0.22), H, GOLD)
txt(s, "⚽  FIFA WORLD CUP 2026",
    Inches(0.4), Inches(0.8), W - Inches(0.6), Inches(0.8),
    size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "INSPIRATIONAL STORIES APP",
    Inches(0.4), Inches(1.65), W - Inches(0.6), Inches(0.6),
    size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
box(s, Inches(1.5), Inches(2.42), W - Inches(3.0), Inches(0.04), GOLD)
txt(s, "User Guide  ·  How to Navigate and Use the App",
    Inches(0.4), Inches(2.55), W - Inches(0.6), Inches(0.4),
    size=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Envisioned by Sameer Avadhani  ·  High Touch Engineer  ·  CX",
    Inches(0.4), Inches(3.2), W - Inches(0.6), Inches(0.35),
    size=12, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "Created and developed using AI",
    Inches(0.4), Inches(3.55), W - Inches(0.6), Inches(0.3),
    size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "FIFA 2026 Inspirational Stories",
    Inches(0.4), H - Inches(0.35), W - Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2  What is this app?
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "What Is This App?",
          "A daily dose of inspiration from FIFA World Cup 2026 players — delivered to your desktop.")

items = [
    ("⚽", "ONE STORY PER DAY",
     "A new inspiring player story appears every day from June 16 – July 11, 2026."),
    ("📄", "SAVE AS WORD OR POWERPOINT",
     "Every story saves as a .docx article or a .pptx slide deck to your local machine."),
    ("🗂", "ONE SHARED FILE OR INDIVIDUAL FILES",
     "Append all stories into one growing document, or save each player as a separate file."),
    ("🖥", "LIVES IN YOUR MENU BAR",
     "The ⚽ icon sits in your Mac menu bar or Windows system tray — always one click away."),
    ("🔔", "OPTIONAL DAILY POPUP",
     "Choose to see today's story automatically every time you open your computer."),
]
y = Inches(1.15)
for icon, heading, detail in items:
    box(s, Inches(0.3), y, Inches(0.42), Inches(0.32), GREEN)
    txt(s, icon, Inches(0.3), y - Inches(0.02), Inches(0.42), Inches(0.36),
        size=14, align=PP_ALIGN.CENTER)
    txt(s, heading, Inches(0.82), y, Inches(2.4), Inches(0.32),
        size=11, bold=True, color=GOLD)
    txt(s, detail,  Inches(3.32), y, Inches(6.4), Inches(0.32),
        size=11, color=WHITE)
    y += Inches(0.72)

footer(s)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3  Menu Bar
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "The ⚽ Menu Bar Icon  ·  Your Starting Point",
          "Click the ⚽ icon in your Mac menu bar (top-right) or Windows system tray (bottom-right).")

menu = [
    (GREEN,  "⚽  Today's Story",
     "Opens today's player story popup immediately."),
    (NAVY,   "📖  Browse All Stories",
     "Opens the player gallery — click any card to read that player's story."),
    (NAVY,   "⚙   Settings",
     "Change your default save folder, file format, and startup preferences."),
    (DARK,   "✖  Quit",
     "Closes the app completely. The icon disappears from the menu bar."),
]
y = Inches(1.15)
for bg_color, label, desc in menu:
    box(s, Inches(0.3), y, Inches(3.5), Inches(0.42), bg_color)
    lc = GOLD if bg_color == GREEN else WHITE
    txt(s, label, Inches(0.38), y + Inches(0.06), Inches(3.35), Inches(0.32),
        size=12, bold=True, color=lc)
    txt(s, desc, Inches(4.0), y + Inches(0.06), Inches(5.7), Inches(0.32),
        size=11, color=WHITE)
    y += Inches(0.62)

footer(s)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4  Story Popup – top bars
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "Story Popup  ·  Save Location & Save Mode Bars",
          "Two control bars appear at the top of every story — giving you full control before saving.")

# Orange bar mock
box(s, Inches(0.3), Inches(1.18), W - Inches(0.6), Inches(0.5), NAVY)
box(s, Inches(0.3), Inches(1.18), Inches(0.06), Inches(0.5), ORANGE)
txt(s, "📁  Save to:   /Users/you/FIFA_Inspirational_Stories",
    Inches(0.45), Inches(1.22), Inches(6.8), Inches(0.38),
    size=11, color=ORANGE)
box(s, Inches(7.4), Inches(1.23), Inches(1.5), Inches(0.38), ORANGE)
txt(s, "Browse…", Inches(7.4), Inches(1.27), Inches(1.5), Inches(0.3),
    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "ORANGE BAR — Save Location",
    Inches(0.3), Inches(1.75), Inches(4.0), Inches(0.28),
    size=12, bold=True, color=ORANGE)
txt(s, "Shows where files will be saved. Type a path or click Browse… to pick any folder on your machine.\nThis overrides the global setting just for this one story.",
    Inches(0.3), Inches(2.06), W - Inches(0.6), Inches(0.44),
    size=11, color=WHITE)

# Blue bar mock
box(s, Inches(0.3), Inches(2.65), W - Inches(0.6), Inches(0.5), NAVY)
box(s, Inches(0.3), Inches(2.65), Inches(0.06), Inches(0.5), BLUE)
txt(s, "💾  Save mode:   ◉ Append to shared FIFA document     ○ Save as individual file for this player",
    Inches(0.45), Inches(2.7), W - Inches(0.65), Inches(0.38),
    size=11, color=BLUE)

txt(s, "BLUE BAR — Save Mode",
    Inches(0.3), Inches(3.22), Inches(4.0), Inches(0.28),
    size=12, bold=True, color=BLUE)

rows = [
    ("◉  Append to shared FIFA document",
     "Adds this story to FIFA_Inspirational_Stories_2026.docx / .pptx — one growing file for all players."),
    ("○  Save as individual file for this player",
     "Creates a standalone file named  2026-06-16_Cristiano_Ronaldo.docx / .pptx — complete with cover + Thank You."),
]
y = Inches(3.55)
for label, detail in rows:
    txt(s, label, Inches(0.3), y, Inches(4.0), Inches(0.28),
        size=11, bold=True, color=WHITE)
    txt(s, detail, Inches(4.4), y, Inches(5.3), Inches(0.32),
        size=11, color=GREY)
    y += Inches(0.52)

footer(s)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5  Story Popup – content sections
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "Story Popup  ·  Reading the Story",
          "Scroll down inside the popup window to read the full story.")

sections = [
    (GREEN,  "APP HEADER",
     "⚽ FIFA WORLD CUP 2026 | Daily Inspiration  — appears at the top of every popup."),
    (DARK,   "PLAYER NAME & INFO",
     "Large gold player name. Below it: Country · Position · Date. Then the story headline in italic gold."),
    (DARK,   "STORY BODY",
     "Full multi-paragraph story. SCROLL DOWN to read all of it. Use the A−  A  A+  A++ buttons (top-right of the window) to adjust text size."),
    (DARK,   "QUOTE BLOCK",
     "The player's most powerful quote, displayed in gold italic between borders."),
    (DARK,   "TODAY'S LESSON",
     "The key takeaway from the player's story — a short lesson in italic text."),
]
y = Inches(1.15)
row_h = Inches(0.58)
for color, heading, detail in sections:
    box(s, Inches(0.3), y, Inches(0.06), row_h, GOLD)
    box(s, Inches(0.36), y, Inches(2.5), row_h, color)
    txt(s, heading, Inches(0.44), y + Inches(0.14), Inches(2.35), Inches(0.3),
        size=10, bold=True, color=GOLD)
    txt(s, detail, Inches(2.96), y + Inches(0.1), Inches(6.7), Inches(row_h - Inches(0.1)),
        size=11, color=WHITE)
    y += row_h + Inches(0.06)

footer(s)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6  Action Buttons
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "Story Popup  ·  The 5 Action Buttons",
          "Five buttons appear at the bottom of every story popup.")

buttons = [
    (GREEN,  WHITE,  "Save as Word (.docx)",
     "Saves the story as a Word document to your chosen folder.\nUses the Save Mode selected in the blue bar above."),
    (GREEN,  WHITE,  "Save as Slide (.pptx)",
     "Saves the story as a PowerPoint slide deck.\nIncludes cover page, full story slides, and a Thank You slide."),
    (GOLD,   DARK,   "Save Both",
     "Saves both .docx and .pptx at once.\nRecommended for a complete archive."),
    (NAVY,   WHITE,  "Open Save Folder",
     "Opens the save folder in Finder (Mac) or File Explorer (Windows)\nso you can immediately see and share your saved files."),
    (RGBColor(0x44,0x44,0x44), WHITE, "Close",
     "Closes this story popup.\nThe ⚽ menu bar icon stays active — the app keeps running."),
]
y = Inches(1.15)
for bg_c, tc, label, desc in buttons:
    box(s, Inches(0.3), y, Inches(2.2), Inches(0.46), bg_c)
    txt(s, label, Inches(0.3), y + Inches(0.08), Inches(2.2), Inches(0.3),
        size=11, bold=True, color=tc, align=PP_ALIGN.CENTER)
    txt(s, desc, Inches(2.65), y + Inches(0.04), Inches(6.9), Inches(0.44),
        size=11, color=WHITE)
    y += Inches(0.64)

footer(s)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7  Browse All Stories
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "Browse All Stories  ·  The Player Gallery",
          "Menu bar ⚽ → Browse All Stories")

txt(s, "The gallery shows all 26 players as clickable cards.\nClick any card to open that player's full story in a new window.",
    Inches(0.3), Inches(1.12), W - Inches(0.6), Inches(0.52),
    size=12, color=WHITE)

# Sample cards  (3 across)
cards = [
    ("Jun 16", "Cristiano Ronaldo", "Portugal · Forward"),
    ("Jun 17", "Vozinha",           "Cabo Verde · GK"),
    ("Jun 18", "Lionel Messi",      "Argentina · Forward"),
    ("Jun 19", "Kylian Mbappé",     "France · Forward"),
    ("Jun 20", "Sadio Mané",        "Senegal · Forward"),
    ("Jun 21", "Luka Modrić",       "Croatia · MF"),
]
cw, ch = Inches(2.85), Inches(0.88)
gap    = Inches(0.2)
for i, (dt, name, meta) in enumerate(cards):
    col = i % 3
    row = i // 3
    cx  = Inches(0.3) + col * (cw + gap)
    cy  = Inches(1.72) + row * (ch + gap)
    box(s, cx, cy, cw, ch, RGBColor(0x22, 0x22, 0x22))
    box(s, cx, cy, cw, Inches(0.06), GREEN)
    txt(s, dt,   cx + Inches(0.1), cy + Inches(0.08), cw - Inches(0.15), Inches(0.22),
        size=9, color=GREY)
    txt(s, name, cx + Inches(0.1), cy + Inches(0.28), cw - Inches(0.15), Inches(0.3),
        size=11, bold=True, color=GOLD)
    txt(s, meta, cx + Inches(0.1), cy + Inches(0.58), cw - Inches(0.15), Inches(0.22),
        size=9, color=GREY)

txt(s, "All 26 players are available here at any time —\nnot just today's story.",
    Inches(9.1) - Inches(0.3), Inches(1.72), Inches(0.0), Inches(0.0),  # placeholder
    size=1, color=DARK)  # invisible spacer

# Right side explanation
rx = Inches(9.1)
# Actually put callout text below the cards
txt(s, "📅  Each card shows the story date, player name, country and position.",
    Inches(0.3), Inches(3.68), W - Inches(0.6), Inches(0.32),
    size=11, color=WHITE)
txt(s, "🖱  Click any card to open the full story popup in a new window — you can have multiple stories open at once.",
    Inches(0.3), Inches(4.02), W - Inches(0.6), Inches(0.32),
    size=11, color=WHITE)
txt(s, "🗓  Stories are ordered by date. Jun 16 = Ronaldo through Jul 11 = Timothy Weah.",
    Inches(0.3), Inches(4.36), W - Inches(0.6), Inches(0.32),
    size=11, color=GOLD)

footer(s)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8  Settings
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "Settings Window",
          "Menu bar ⚽ → Settings")

settings = [
    ("📁", "Save Directory",
     "The default folder where all saved files are stored.\nClick Browse… to change it to any folder on your machine."),
    ("📄", "Output Format",
     "Choose what to save:\n  • Both (.docx + .pptx)  — recommended\n  • Word only (.docx)\n  • Slide only (.pptx)"),
    ("🔁", "Launch at Login",
     "When checked, the ⚽ icon appears automatically in your menu bar\nevery time you log in to your computer."),
    ("🔔", "Show Story at Startup",
     "When checked, today's player story popup opens automatically\nwhen you log in. Uncheck to open it manually."),
]
y = Inches(1.18)
for icon, heading, detail in settings:
    box(s, Inches(0.3), y, Inches(0.5), Inches(0.5), GREEN)
    txt(s, icon,    Inches(0.3), y + Inches(0.08), Inches(0.5), Inches(0.36),
        size=16, align=PP_ALIGN.CENTER)
    txt(s, heading, Inches(0.9), y + Inches(0.02), Inches(2.2), Inches(0.28),
        size=12, bold=True, color=GOLD)
    txt(s, detail,  Inches(3.2), y, Inches(6.4), Inches(0.55),
        size=11, color=WHITE)
    y += Inches(0.78)

footer(s)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9  Saved Files
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "Your Saved Files",
          "All files save to your local machine only — nothing is uploaded or shared.")

files = [
    ("📘", "FIFA_Inspirational_Stories_2026.docx",
     "SHARED (Append mode) — one Word document. Stories are added as new pages each time you save.\nIncludes a cover page with credits."),
    ("📊", "FIFA_Inspirational_Stories_2026.pptx",
     "SHARED (Append mode) — one PowerPoint deck. Stories are added as new slides.\nIncludes cover + Thank You slide."),
    ("👤", "2026-06-16_Cristiano_Ronaldo.docx / .pptx",
     "INDIVIDUAL (Individual mode) — a standalone file per player.\nContains cover, story, and Thank You slide. Complete and self-contained."),
    ("🔒", "No duplicate stories",
     "The app tracks what has already been saved. Clicking Save again on the same story\nin Append mode does nothing — it won't create duplicate pages."),
]
y = Inches(1.15)
for icon, filename, detail in files:
    box(s, Inches(0.3), y, Inches(0.5), Inches(0.9), GREEN)
    txt(s, icon, Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.5),
        size=18, align=PP_ALIGN.CENTER)
    txt(s, filename, Inches(0.9), y + Inches(0.04), Inches(8.6), Inches(0.3),
        size=11, bold=True, color=GOLD)
    txt(s, detail,   Inches(0.9), y + Inches(0.36), Inches(8.6), Inches(0.52),
        size=10, color=WHITE)
    y += Inches(1.0)

footer(s)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10  Story Schedule
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "Daily Story Schedule  ·  June 16 – July 11, 2026",
          "One player story per day — automatically shown based on today's date.")

schedule = [
    ("Jun 16", "Cristiano Ronaldo",   "Portugal"),
    ("Jun 17", "Vozinha (J. Dias)",   "Cabo Verde"),
    ("Jun 18", "Lionel Messi",        "Argentina"),
    ("Jun 19", "Kylian Mbappé",       "France"),
    ("Jun 20", "Sadio Mané",          "Senegal"),
    ("Jun 21", "Luka Modrić",         "Croatia"),
    ("Jun 22", "Mohamed Salah",       "Egypt"),
    ("Jun 23", "Son Heung-min",       "South Korea"),
    ("Jun 24", "N'Golo Kanté",        "France"),
    ("Jun 25", "Vinicius Jr.",        "Brazil"),
    ("Jun 26", "R. Lewandowski",      "Poland"),
    ("Jun 27", "Hakim Ziyech",        "Morocco"),
    ("Jun 28", "Achraf Hakimi",       "Morocco"),
    ("Jun 29", "Keylor Navas",        "Costa Rica"),
    ("Jun 30", "Marcus Rashford",     "England"),
    ("Jul 01", "Riyad Mahrez",        "Algeria"),
    ("Jul 02", "Pedri",               "Spain"),
    ("Jul 03", "Jude Bellingham",     "England"),
    ("Jul 04", "Gavi",                "Spain"),
    ("Jul 05", "Bukayo Saka",         "England"),
    ("Jul 06", "Rodri",               "Spain"),
    ("Jul 07", "Erling Haaland",      "Norway"),
    ("Jul 08", "Camavinga",           "France"),
    ("Jul 09", "Lamine Yamal",        "Spain"),
    ("Jul 10", "Andrés Guardado",     "Mexico"),
    ("Jul 11", "Timothy Weah",        "USA"),
]

cols   = 4
col_w  = (W - Inches(0.6)) / cols
row_h  = (H - Inches(1.45)) / 7   # 7 rows for 26 items

for i, (dt, name, country) in enumerate(schedule):
    col = i % cols
    row = i // cols
    x   = Inches(0.3) + col * col_w
    y   = Inches(1.12) + row * row_h
    bg_c = GREEN if col % 2 == 0 else RGBColor(0x1e, 0x1e, 0x1e)
    box(s, x, y, col_w - Inches(0.08), row_h - Inches(0.04), bg_c)
    txt(s, dt,
        x + Inches(0.06), y + Inches(0.01),
        col_w - Inches(0.14), row_h * 0.38,
        size=8, bold=True, color=GOLD)
    txt(s, name,
        x + Inches(0.06), y + row_h * 0.36,
        col_w - Inches(0.14), row_h * 0.38,
        size=8, bold=False, color=WHITE)
    txt(s, country,
        x + Inches(0.06), y + row_h * 0.68,
        col_w - Inches(0.14), row_h * 0.3,
        size=7, color=GREY)

footer(s)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11  Thank You
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, GREEN)
box(s, 0, 0, Inches(0.22), H, GOLD)
txt(s, "Thank You",
    Inches(0.4), Inches(1.4), W - Inches(0.6), Inches(0.9),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "for taking a moment to be inspired.",
    Inches(0.4), Inches(2.4), W - Inches(0.6), Inches(0.5),
    size=18, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
box(s, Inches(1.5), Inches(3.1), W - Inches(3.0), Inches(0.04), GOLD)
txt(s, "Envisioned by Sameer Avadhani  ·  High Touch Engineer  ·  CX",
    Inches(0.4), Inches(3.25), W - Inches(0.6), Inches(0.35),
    size=12, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "Created and developed using AI",
    Inches(0.4), Inches(3.62), W - Inches(0.6), Inches(0.3),
    size=11, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "FIFA 2026 Inspirational Stories",
    Inches(0.4), H - Inches(0.35), W - Inches(0.6), Inches(0.3),
    size=9, color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)


prs.save(str(OUT))
print(f"Saved: {OUT}  ({OUT.stat().st_size:,} bytes)")
