"""
generator.py
- One persistent FIFA_Inspirational_Stories.docx  (cover page + stories appended)
- One persistent FIFA_Inspirational_Stories.pptx  (title slide + one slide per story)
- Stories are only appended once; re-saving the same story is a no-op.
"""
from pathlib import Path
from datetime import date, datetime

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN

import json, os

# ── File names (single persistent files) ──────────────────────────────────────
DOCX_NAME  = "FIFA_Inspirational_Stories_2026.docx"
PPTX_NAME  = "FIFA_Inspirational_Stories_2026.pptx"
INDEX_NAME = ".published_stories.json"          # hidden tracker

# ── Colors ─────────────────────────────────────────────────────────────────────
GOLD   = RGBColor(0xD4, 0xAF, 0x37)
GREEN  = RGBColor(0x00, 0x6B, 0x3C)
BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x88, 0x88, 0x88)

PGOLD  = PptxRGB(0xD4, 0xAF, 0x37)
PGREEN = PptxRGB(0x00, 0x6B, 0x3C)
PWHITE = PptxRGB(0xFF, 0xFF, 0xFF)
PDARK  = PptxRGB(0x1A, 0x1A, 0x1A)
PGREY  = PptxRGB(0x88, 0x88, 0x88)

CREDITS  = "Envisioned by Sameer Avadhani · High Touch Engineer · CX\nCreated and developed using AI"
FOOTER_TEXT = "FIFA 2026 Inspirational Stories"


def _ensure_dir(directory: str) -> Path:
    p = Path(directory)
    p.mkdir(parents=True, exist_ok=True)
    return p


# ── Published-story tracker ────────────────────────────────────────────────────
def _load_index(directory: str) -> dict:
    f = Path(directory) / INDEX_NAME
    if f.exists():
        try:
            return json.loads(f.read_text())
        except Exception:
            pass
    return {"docx": [], "pptx": []}


def _save_index(directory: str, index: dict):
    f = Path(directory) / INDEX_NAME
    f.write_text(json.dumps(index, indent=2))


def _already_published(directory: str, story_id: int, fmt: str) -> bool:
    return story_id in _load_index(directory).get(fmt, [])


def _mark_published(directory: str, story_id: int, fmt: str):
    idx = _load_index(directory)
    if story_id not in idx.get(fmt, []):
        idx.setdefault(fmt, []).append(story_id)
    _save_index(directory, idx)


# ── WORD helpers ───────────────────────────────────────────────────────────────
def _set_cell_bg(cell, hex_color: str):
    """Fill a table cell background color."""
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement("w:shd")
    shd.set(qn("w:val"),   "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"),  hex_color)
    tcPr.append(shd)


def _add_run(para, text, size=12, bold=False, italic=False,
             color: RGBColor = None):
    r = para.add_run(text)
    r.bold        = bold
    r.italic      = italic
    r.font.size   = Pt(size)
    if color:
        r.font.color.rgb = color
    return r


def _para(doc, align=WD_ALIGN_PARAGRAPH.LEFT, space_before=0, space_after=6):
    p = doc.add_paragraph()
    p.alignment         = align
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after  = Pt(space_after)
    return p


def _page_break(doc):
    doc.add_page_break()


# ── WORD cover page ────────────────────────────────────────────────────────────
def _build_word_cover(doc: Document):
    # Green banner table as "cover"
    tbl = doc.add_table(rows=1, cols=1)
    tbl.style = "Table Grid"
    cell = tbl.cell(0, 0)
    _set_cell_bg(cell, "006B3C")
    cell.width = Inches(6.5)

    cp = cell.paragraphs[0]
    cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cp.paragraph_format.space_before = Pt(48)
    cp.paragraph_format.space_after  = Pt(8)
    _add_run(cp, "⚽  FIFA WORLD CUP 2026", 28, bold=True, color=WHITE)

    cp2 = cell.add_paragraph()
    cp2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cp2.paragraph_format.space_after = Pt(8)
    _add_run(cp2, "INSPIRATIONAL STORIES", 22, bold=True, color=GOLD)

    cp3 = cell.add_paragraph()
    cp3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cp3.paragraph_format.space_after = Pt(48)
    _add_run(cp3, "One Player. One Story. Every Day.", 13, italic=True, color=WHITE)

    doc.add_paragraph()

    # Credits box
    cb = doc.add_paragraph()
    cb.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _add_run(cb, CREDITS.replace("\n", "  ·  "), 11, italic=True, color=GREEN)

    doc.add_paragraph()

    ts = doc.add_paragraph()
    ts.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _add_run(ts, f"Published: June 16 – July 11, 2026", 11, color=GREY)

    _page_break(doc)


# ── WORD story page ────────────────────────────────────────────────────────────
def _append_word_story(doc: Document, story: dict):
    # Player name header
    h = _para(doc, WD_ALIGN_PARAGRAPH.CENTER, space_after=2)
    _add_run(h, story["name"].upper(), 26, bold=True, color=GOLD)

    meta = _para(doc, WD_ALIGN_PARAGRAPH.CENTER, space_after=4)
    _add_run(meta, f"{story['country']}  ·  {story['position']}  ·  {story['date']}",
             11, color=GREY)

    hl = _para(doc, WD_ALIGN_PARAGRAPH.CENTER, space_after=10)
    _add_run(hl, f'"{story["headline"]}"', 14, bold=True, italic=True, color=GREEN)

    # Divider
    div = _para(doc, WD_ALIGN_PARAGRAPH.CENTER, space_after=8)
    _add_run(div, "─" * 60, 9, color=GREY)

    # Story body — split paragraphs on \n\n
    for block in story["story"].split("\n\n"):
        bp = _para(doc, WD_ALIGN_PARAGRAPH.JUSTIFY, space_after=8)
        _add_run(bp, block.strip(), 12)

    doc.add_paragraph()

    # Quote
    qt = _para(doc, WD_ALIGN_PARAGRAPH.CENTER, space_after=8)
    _add_run(qt, f'❝  {story["quote"]}  ❞', 13, bold=True, italic=True, color=GOLD)

    # Lesson
    ls_row = doc.add_table(rows=1, cols=2)
    ls_row.style = "Table Grid"
    _set_cell_bg(ls_row.cell(0, 0), "006B3C")
    _set_cell_bg(ls_row.cell(0, 1), "1A1A1A")

    lbl = ls_row.cell(0, 0).paragraphs[0]
    lbl.alignment = WD_ALIGN_PARAGRAPH.CENTER
    lbl.paragraph_format.space_before = Pt(6)
    lbl.paragraph_format.space_after  = Pt(6)
    _add_run(lbl, "TODAY'S\nLESSON", 10, bold=True, color=WHITE)

    txt = ls_row.cell(0, 1).paragraphs[0]
    txt.paragraph_format.space_before = Pt(4)
    txt.paragraph_format.space_after  = Pt(4)
    _add_run(txt, story["lesson"], 11, italic=True, color=WHITE)

    doc.add_paragraph()
    _page_break(doc)


# ── PPTX helpers ───────────────────────────────────────────────────────────────
def _new_pptx() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Emu(9144000)   # 10 in
    prs.slide_height = Emu(5143500)   # 7.5 in
    return prs


def _pptx_rect(slide, left, top, w, h, fill):
    from pptx.util import Emu
    sh = slide.shapes.add_shape(1, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def _pptx_tb(slide, text, left, top, w, h, size, bold=False, italic=False,
             color=PWHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text        = text
    run.font.size   = PPt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color


# ── PPTX title slide ───────────────────────────────────────────────────────────
def _build_pptx_cover(prs: Presentation):
    W, H   = prs.slide_width, prs.slide_height
    slide  = prs.slides.add_slide(prs.slide_layouts[6])

    _pptx_rect(slide, 0, 0, W, H, PGREEN)
    _pptx_rect(slide, 0, 0, Emu(28000), H, PGOLD)        # gold left bar

    _pptx_tb(slide, "⚽  FIFA WORLD CUP 2026",
             PInches(0.4), PInches(1.2), W - PInches(0.8), PInches(0.9),
             size=44, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)

    _pptx_tb(slide, "INSPIRATIONAL STORIES",
             PInches(0.4), PInches(2.1), W - PInches(0.8), PInches(0.7),
             size=32, bold=True, color=PGOLD, align=PP_ALIGN.CENTER)

    _pptx_tb(slide, "One Player. One Story. Every Day.",
             PInches(0.4), PInches(2.85), W - PInches(0.8), PInches(0.4),
             size=16, italic=True, color=PWHITE, align=PP_ALIGN.CENTER)

    _pptx_rect(slide, PInches(1.5), PInches(3.4), W - PInches(3.0), Emu(18000), PGOLD)

    _pptx_tb(slide, CREDITS,
             PInches(0.4), PInches(3.55), W - PInches(0.8), PInches(0.7),
             size=12, italic=True, color=PGOLD, align=PP_ALIGN.CENTER)

    _pptx_tb(slide, "June 16 – July 11, 2026",
             PInches(0.4), PInches(4.3), W - PInches(0.8), PInches(0.35),
             size=11, color=PGREY, align=PP_ALIGN.CENTER)


# ── PPTX story slides (multi-slide: header + body paragraphs + closing) ────────
def _append_pptx_story(prs: Presentation, story: dict):
    W, H = prs.slide_width, prs.slide_height

    def footer(slide):
        _pptx_tb(slide, FOOTER_TEXT,
                 PInches(0.35), H - PInches(0.26), W - PInches(0.5), PInches(0.22),
                 size=8, color=PGREY, align=PP_ALIGN.CENTER)

    # ── Slide 1: Header ──────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_rect(s1, 0, 0, W, H, PDARK)
    _pptx_rect(s1, 0, 0, Emu(18000), H, PGREEN)

    _pptx_tb(s1, "⚽  FIFA WORLD CUP 2026  |  Daily Inspiration",
             PInches(0.35), PInches(0.18), W - PInches(0.5), PInches(0.38),
             size=13, bold=True, color=PGOLD, align=PP_ALIGN.CENTER)
    _pptx_rect(s1, PInches(0.35), PInches(0.62), W - PInches(0.7), Emu(22000), PGOLD)

    _pptx_tb(s1, story["name"].upper(),
             PInches(0.35), PInches(0.7), W - PInches(0.5), PInches(0.72),
             size=38, bold=True, color=PGOLD, align=PP_ALIGN.CENTER)
    _pptx_tb(s1, f"{story['country']}  ·  {story['position']}  ·  {story['date']}",
             PInches(0.35), PInches(1.45), W - PInches(0.5), PInches(0.35),
             size=13, color=PWHITE, align=PP_ALIGN.CENTER)
    _pptx_tb(s1, f'"{story["headline"]}"',
             PInches(0.5), PInches(1.88), W - PInches(1.0), PInches(0.55),
             size=17, bold=True, italic=True, color=PGOLD, align=PP_ALIGN.CENTER)
    _pptx_rect(s1, PInches(0.35), PInches(2.55), W - PInches(0.7), Emu(16000), PGREEN)

    # First paragraph of story on slide 1
    paragraphs = [p.strip() for p in story["story"].split("\n\n") if p.strip()]
    first_para = paragraphs[0] if paragraphs else ""
    _pptx_tb(s1, first_para,
             PInches(0.45), PInches(2.65), W - PInches(0.9), PInches(1.95),
             size=11, color=PWHITE, align=PP_ALIGN.JUSTIFY)
    footer(s1)

    # ── Slides 2..N: Remaining paragraphs (2 per slide) ─────────────────────
    remaining = paragraphs[1:]
    paired    = [remaining[i:i+2] for i in range(0, len(remaining), 2)]

    for pair in paired:
        sx = prs.slides.add_slide(prs.slide_layouts[6])
        _pptx_rect(sx, 0, 0, W, H, PDARK)
        _pptx_rect(sx, 0, 0, Emu(18000), H, PGREEN)
        _pptx_tb(sx, story["name"].upper() + "  (cont.)",
                 PInches(0.35), PInches(0.12), W - PInches(0.5), PInches(0.35),
                 size=12, bold=True, color=PGOLD, align=PP_ALIGN.LEFT)
        _pptx_rect(sx, PInches(0.35), PInches(0.52), W - PInches(0.7), Emu(14000), PGREEN)
        combined = "\n\n".join(pair)
        _pptx_tb(sx, combined,
                 PInches(0.45), PInches(0.62), W - PInches(0.9), PInches(4.3),
                 size=12, color=PWHITE, align=PP_ALIGN.JUSTIFY)
        footer(sx)

    # ── Final slide: Quote + Lesson ──────────────────────────────────────────
    sf = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_rect(sf, 0, 0, W, H, PDARK)
    _pptx_rect(sf, 0, 0, Emu(18000), H, PGREEN)
    _pptx_rect(sf, PInches(0.35), PInches(0.15), W - PInches(0.7), Emu(16000), PGOLD)
    _pptx_tb(sf, story["name"].upper() + "  –  Key Takeaways",
             PInches(0.35), PInches(0.22), W - PInches(0.5), PInches(0.35),
             size=14, bold=True, color=PDARK, align=PP_ALIGN.CENTER)

    _pptx_tb(sf, f'❝  {story["quote"]}  ❞',
             PInches(0.5), PInches(0.75), W - PInches(1.0), PInches(1.0),
             size=16, bold=True, italic=True, color=PGOLD, align=PP_ALIGN.CENTER)
    _pptx_rect(sf, PInches(0.35), PInches(1.85), W - PInches(0.7), Emu(16000), PGREEN)
    _pptx_tb(sf, "TODAY'S LESSON",
             PInches(0.45), PInches(1.92), PInches(2.5), PInches(0.3),
             size=11, bold=True, color=PGOLD)
    _pptx_tb(sf, story["lesson"],
             PInches(0.45), PInches(2.28), W - PInches(0.9), PInches(1.8),
             size=12, italic=True, color=PWHITE, align=PP_ALIGN.JUSTIFY)
    footer(sf)


# ── PPTX entry point ───────────────────────────────────────────────────────────
def generate_pptx(story: dict, directory: str, individual: bool = False) -> str:
    out_dir = _ensure_dir(directory)

    if individual:
        # Standalone file per player — always regenerate fresh
        safe  = story["name"].replace(" ", "_").replace("(", "").replace(")", "")
        fpath = out_dir / f"{story['date']}_{safe}.pptx"
        prs   = _new_pptx()
        _build_pptx_cover(prs)
        _append_pptx_story(prs, story)
        _build_pptx_thankyou(prs)
        prs.save(str(fpath))
        return str(fpath)

    # Shared appended file
    filepath = out_dir / PPTX_NAME
    if _already_published(directory, story["id"], "pptx"):
        return str(filepath)
    prs = Presentation(str(filepath)) if filepath.exists() else _new_pptx()
    if not filepath.exists():
        _build_pptx_cover(prs)
    _append_pptx_story(prs, story)
    prs.save(str(filepath))
    _mark_published(directory, story["id"], "pptx")
    return str(filepath)


# ── PPTX Thank You slide ───────────────────────────────────────────────────────
def _build_pptx_thankyou(prs: Presentation):
    W, H  = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_rect(slide, 0, 0, W, H, PGREEN)
    _pptx_rect(slide, 0, 0, Emu(18000), H, PGOLD)
    _pptx_tb(slide, "Thank You",
             PInches(0.4), PInches(1.6), W - PInches(0.8), PInches(1.0),
             size=54, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)
    _pptx_tb(slide, "for taking a moment to be inspired.",
             PInches(0.4), PInches(2.7), W - PInches(0.8), PInches(0.5),
             size=18, italic=True, color=PGOLD, align=PP_ALIGN.CENTER)
    _pptx_rect(slide, PInches(1.5), PInches(3.4), W - PInches(3.0), Emu(16000), PGOLD)
    _pptx_tb(slide, CREDITS,
             PInches(0.4), PInches(3.55), W - PInches(0.8), PInches(0.6),
             size=11, italic=True, color=PWHITE, align=PP_ALIGN.CENTER)
    _pptx_tb(slide, FOOTER_TEXT,
             PInches(0.4), H - PInches(0.3), W - PInches(0.8), PInches(0.25),
             size=9, color=PGREY, align=PP_ALIGN.CENTER)


# ── WORD Thank You page ────────────────────────────────────────────────────────
def _append_word_thankyou(doc: Document):
    _page_break(doc)
    tbl = doc.add_table(rows=1, cols=1)
    tbl.style = "Table Grid"
    cell = tbl.cell(0, 0)
    _set_cell_bg(cell, "006B3C")
    p1 = cell.paragraphs[0]
    p1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p1.paragraph_format.space_before = Pt(60)
    p1.paragraph_format.space_after  = Pt(8)
    _add_run(p1, "Thank You", 36, bold=True, color=WHITE)
    p2 = cell.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2.paragraph_format.space_after = Pt(8)
    _add_run(p2, "for taking a moment to be inspired.", 14, italic=True, color=GOLD)
    p3 = cell.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p3.paragraph_format.space_after = Pt(60)
    _add_run(p3, CREDITS.replace("\n", "  ·  "), 10, italic=True, color=WHITE)
    doc.add_paragraph()
    ft = _para(doc, WD_ALIGN_PARAGRAPH.CENTER)
    _add_run(ft, FOOTER_TEXT, 10, color=GREY)


# ── WORD entry point ───────────────────────────────────────────────────────────
def generate_word(story: dict, directory: str, individual: bool = False) -> str:
    out_dir = _ensure_dir(directory)

    if individual:
        safe  = story["name"].replace(" ", "_").replace("(", "").replace(")", "")
        fpath = out_dir / f"{story['date']}_{safe}.docx"
        doc   = Document()
        for section in doc.sections:
            section.top_margin    = Inches(0.8)
            section.bottom_margin = Inches(0.8)
            section.left_margin   = Inches(1.0)
            section.right_margin  = Inches(1.0)
        _build_word_cover(doc)
        _append_word_story(doc, story)
        _append_word_thankyou(doc)
        doc.save(str(fpath))
        return str(fpath)

    # Shared appended file
    filepath = out_dir / DOCX_NAME
    if _already_published(directory, story["id"], "docx"):
        return str(filepath)
    if filepath.exists():
        doc = Document(str(filepath))
    else:
        doc = Document()
        for section in doc.sections:
            section.top_margin    = Inches(0.8)
            section.bottom_margin = Inches(0.8)
            section.left_margin   = Inches(1.0)
            section.right_margin  = Inches(1.0)
        _build_word_cover(doc)
    _append_word_story(doc, story)
    doc.save(str(filepath))
    _mark_published(directory, story["id"], "docx")
    return str(filepath)


# ── Public API ─────────────────────────────────────────────────────────────────
def generate_documents(story: dict, directory: str, fmt: str = "both",
                       individual: bool = False) -> list[str]:
    paths = []
    if fmt in ("docx", "both"):
        paths.append(generate_word(story, directory, individual))
    if fmt in ("pptx", "both"):
        paths.append(generate_pptx(story, directory, individual))
    return paths
